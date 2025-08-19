import calendar
import json
import math
import pdb

import pandas as pd
import plotly.express as px
import plotly.graph_objs as go
import plotly.io as pio
import requests

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR

import datetime

# Import project configuration instead of company-specific values
from config import Device_to_location, compliance_data_link


# Functions

def get_compliance_handwash(x):
    """Return 1 if the handwash observed lasted at least 40 seconds."""
    if x['Hand Washes Observed']:
        if x['TotalWashDuration'] >= 40:
            return 1
    return 0


def get_100_compliance_handwash(x):
    """Return 1 if the handwash lasted at least 40s and all core steps were performed."""
    if x['Hand Washes with Total Duration >= 40 sec']:
        if x['HandWash With All Core Steps']:
            return 1
    return 0


def text_into_presentation(pres, slide_no, shape_no, para_no, run_no, words):
    """Replace text in a given placeholder of a PowerPoint slide."""
    try:
        pres.slides[slide_no].shapes[shape_no].text_frame.paragraphs[para_no].runs[run_no].text = words
        print(f"Updated text on slide {slide_no}, shape {shape_no}")
    except Exception as e:
        print(f"Failed to update text: {e}")


def table_cell_into_presentation(pres, slide_no, shape_no, cell_x, cell_y, words):
    """Insert text into a cell of a PowerPoint table."""
    try:
        pres.slides[slide_no].shapes[shape_no].table.cell(cell_x, cell_y).text = words
        print(f"Updated table cell ({cell_x}, {cell_y}) on slide {slide_no}")
    except Exception as e:
        print(f"Failed to update table cell: {e}")


def img_into_presentation(pres, slide_no, path, pos_x, pos_y, size_x, size_y):
    """Insert an image into the presentation at a given location."""
    try:
        pres.slides[slide_no].shapes.add_picture(
            path,
            Inches(pos_x),
            Inches(pos_y),
            width=Inches(size_x),
            height=Inches(size_y)
        )
        print(f"Inserted image {path} on slide {slide_no}")
    except Exception as e:
        print(f"Failed to insert image: {e}")


def add_text_to_shape(pr_ppt, slide_no, shape_no, text, color, most_common, left, top, width, height):
    """Insert a textbox with styled text into a presentation slide."""
    slide = pr_ppt.slides[slide_no]

    new_shape = slide.shapes.add_textbox(left, top, width, height)
    new_shape.text_frame.text = text
    new_shape.text_frame.word_wrap = True
    new_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)

    # Apply styling depending on "most_common" flag
    for paragraph in new_shape.text_frame.paragraphs:
        paragraph.font.name = 'Poppins Light'
        paragraph.font.size = Pt(11 if most_common else 8)

    if most_common:
        new_shape.text_frame.paragraphs[0].font.bold = True


def add_text_to_shape2(pr_ppt, slide_no, text, color, left, top, width, height, rotation):
    """Add rotated text to a PowerPoint slide."""
    slide = pr_ppt.slides[slide_no]

    new_shape = slide.shapes.add_textbox(left, top, width, height)
    new_shape.text_frame.text = text
    new_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*color)
    new_shape.text_frame.paragraphs[0].font.name = 'Poppins Light'
    new_shape.text_frame.paragraphs[0].font.size = Pt(8)
    new_shape.rotation = rotation


def largest_remainder_rounding(numbers, target_sum):
    """
    Rounds numbers proportionally so that their sum equals the target_sum.
    Uses the largest remainder method.
    """
    sorted_numbers = sorted(numbers, key=lambda x: x % 1, reverse=True)

    floor_sum = sum(math.floor(n) for n in numbers)
    difference = target_sum - floor_sum

    for i in range(difference):
        sorted_numbers[i] = int(sorted_numbers[i]) + 1

    return sorted_numbers
