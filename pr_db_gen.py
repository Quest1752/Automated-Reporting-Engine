import calendar
import json
import numpy

import pandas as pd
import plotly.express as px
import plotly.graph_objs as go
import plotly.io as pio
import requests
# pip3 install kaleido - to use write_image with kaleido engine. 

from pptx import Presentation
from pptx.util import Pt, Inches

from pptx.dml.color import RGBColor

import datetime
import argparse

from pr_values import Device_to_location, compliance_data_link
from pr_gen_functions import text_into_presentation, table_cell_into_presentation, img_into_presentation, \
    add_text_to_shape, add_text_to_shape2

import pdb

# Contains Device_to_location dict

parser = argparse.ArgumentParser()
parser.add_argument('--uuid', default="6a5a8352d7e111eb91880a4165e7dae6")
parser.add_argument('--date_start', default='2023-04-01')
parser.add_argument('--date_end', default='2023-04-30')
parser.add_argument('--time_start', default='00:00:01')
parser.add_argument('--time_end', default='23:59:59')
parser.add_argument('--prev_start', default='2023-03-01')
parser.add_argument('--prev_end', default='2023-03-31')
args = parser.parse_args()

# Variables
print("Setting up constants")

DeviceUUID = args.uuid
start_date = args.date_start
end_date = args.date_end

start_time, end_time = args.time_start, args.time_end
prev_start = args.prev_start
prev_end = args.prev_end

# Extract the day and month from the start and end dates
start_year, start_month, start_day = start_date.split('-')
_, end_month, end_day = end_date.split('-')

# Convert the start and end months from integer to string format
start_month_name = calendar.month_name[int(start_month)]
end_month_name = calendar.month_name[int(end_month)]

date_range_text = f'{start_day} to {end_day} {start_month_name[:3]} {start_year}' #\n{start_time[:5]} to {end_time[:5]}'
date_title_text = f'{end_day} {start_month_name[:3]} {start_year}'
print(date_range_text)
print(date_title_text)

with open('f_factor.json', 'r') as a:
    f_factor = json.load(a)
a.close()

print(f_factor)

# Retrieving data from database
url = "{}/get_specific_compliance".format(compliance_data_link)
# print("current compliance data")

param = {'installations': DeviceUUID, 'start_date': start_date, 'end_date': end_date}
response = requests.get(url, params=param)  # auth=(username, password)) #response = requests.get(url) #fetch_data(url)
d = json.loads(response.text)
df = pd.DataFrame(d)
#pdb.set_trace()
# Making changes to columns
# Current Fortnight
# print("Making changes to columns regarding data types")
df.replace({'True': 1, 'False': 0}, inplace=True)
df['Recorded Videos'] = 1
df['TotalWashDuration'] = df['TotalWashDuration'].astype(float)

df['core_steps_count'] = df[
    ['PalmtoPalm_Detected', 'PalmtoDorsum_Detected', 'FingersInterlaced_Detected', 'FistInterlocked_Detected',
     'ThumbRub_Detected', 'PalmtoNails_Detected']].sum(axis=1)
df['Hand Washes Observed'] = df['core_steps_count'].apply(lambda x: 1 if x > 0 else 0)
df['Low Score'] = df['core_steps_count'].apply(lambda x: 1 if x in [1, 2] else 0)
df['Good Score'] = df['core_steps_count'].apply(lambda x: 1 if x in [3, 4] else 0)
df['High Score'] = df['core_steps_count'].apply(lambda x: 1 if x in [5, 6] else 0)
df['HandWash With All Core Steps'] = df['core_steps_count'].apply(lambda x: 1 if x == 6 else 0)

scaling_factor = f_factor[DeviceUUID]

# df['Hand Washes with Total Duration >= 40 sec'] = df.apply(get_compliance_handwash, axis=1)
# df['Handwash(100% Compliance)'] = df.apply(get_100_compliance_handwash, axis=1)

df['EpisodeTime'] = df['EpisodeTime'].apply(
    lambda x: (str(int(int(x) / 10000)) + ":" + str(int(int(int(x) / 100) % 100))) + ':' + str(int(int(x) % 100)))
df["DateTime"] = df['EpisodeDate'] + " " + df['EpisodeTime']
df["DateTime"] = pd.to_datetime(df["DateTime"])
df['EpisodeDate'] = pd.to_datetime(df['EpisodeDate'])


#df['EpisodeDate'] = pd.to_datetime(df['DateTime'])

#df = df.set_index('DateTime').between_time(start_time, end_time)
#print(df[['EpisodeTime', 'core_steps_count']])
#pdb.set_trace()

low_counts_df = df['Low Score'].sum()  # .sort_index().reset_index()
good_counts_df = df['Good Score'].sum()  # .sort_index().reset_index()
high_counts_df = df['High Score'].sum()  # .sort_index().reset_index()

max_counts = max(low_counts_df, good_counts_df, high_counts_df)

if max_counts == high_counts_df:
    highest_category = 2
elif max_counts == good_counts_df:
    highest_category = 1
elif max_counts == low_counts_df:
    highest_category = 0


#Previous dates values

param = {'installations': DeviceUUID, 'start_date': start_date, 'end_date': end_date}
response = requests.get(url, params=param)  # auth=(username, password)) #response = requests.get(url) #fetch_data(url)
d = json.loads(response.text)
prev_df = pd.DataFrame(d)
#pdb.set_trace()
# Making changes to columns
# Current Fortnight
# print("Making changes to columns regarding data types")
prev_df.replace({'True': 1, 'False': 0}, inplace=True)
prev_df['Recorded Videos'] = 1
prev_df['TotalWashDuration'] = prev_df['TotalWashDuration'].astype(float)

prev_df['core_steps_count'] = prev_df[
    ['PalmtoPalm_Detected', 'PalmtoDorsum_Detected', 'FingersInterlaced_Detected', 'FistInterlocked_Detected',
     'ThumbRub_Detected', 'PalmtoNails_Detected']].sum(axis=1)
prev_df['Hand Washes Observed'] = prev_df['core_steps_count'].apply(lambda x: 1 if x > 0 else 0)
prev_df['Low Score'] = prev_df['core_steps_count'].apply(lambda x: 1 if x in [1, 2] else 0)
prev_df['Good Score'] = prev_df['core_steps_count'].apply(lambda x: 1 if x in [3, 4] else 0)
prev_df['High Score'] = prev_df['core_steps_count'].apply(lambda x: 1 if x in [5, 6] else 0)
prev_df['HandWash With All Core Steps'] = prev_df['core_steps_count'].apply(lambda x: 1 if x == 6 else 0)

scaling_factor = f_factor[DeviceUUID]

# df['Hand Washes with Total Duration >= 40 sec'] = df.apply(get_compliance_handwash, axis=1)
# df['Handwash(100% Compliance)'] = prev_df.apply(get_100_compliance_handwash, axis=1)

prev_df['EpisodeTime'] = prev_df['EpisodeTime'].apply(
    lambda x: (str(int(int(x) / 10000)) + ":" + str(int(int(int(x) / 100) % 100))) + ':' + str(int(int(x) % 100)))
prev_df["DateTime"] = prev_df['EpisodeDate'] + " " + prev_df['EpisodeTime']
prev_df["DateTime"] = pd.to_datetime(prev_df["DateTime"])
# prev_df['EpisodeDate'] = pd.to_datetime(prev_df[['EpisodeDate','EpisodeTime']])


#prev_df['EpisodeDate'] = pd.to_datetime(prev_df['DateTime'])

#prev_df = prev_df.set_index('DateTime').between_time(start_time, end_time)
#print(prev_df[['EpisodeTime', 'core_steps_count']])

# pdb.set_trace()

low_counts_prev_df = prev_df['Low Score'].sum()  # .sort_index().reset_index()
good_counts_prev_df = prev_df['Good Score'].sum()  # .sort_index().reset_index()
high_counts_prev_df = prev_df['High Score'].sum()  # .sort_index().reset_index()

max_counts = max(low_counts_prev_df, good_counts_prev_df, high_counts_prev_df)

if max_counts == high_counts_prev_df:
    prev_category = 2
elif max_counts == good_counts_prev_df:
    prev_category = 1
elif max_counts == low_counts_prev_df:
    prev_category = 0

# Calculate the scaling factor for the current device UUID
#scaling_factor = f_factor[DeviceUUID]
site_type = Device_to_location[DeviceUUID].split(" ")[-1]

# Group the data by date and calculate the sum of relevant columns
datewise_df = df.groupby(['EpisodeDate']).sum().reset_index()

#pdb.set_trace()

# Multiply the 'Hand Washes Observed' column by the scaling factor
datewise_df['Hand Washes Observed'] = datewise_df['Hand Washes Observed'] * scaling_factor

# Convert start_date and end_date to datetime objects
start_date = pd.to_datetime(start_date)
end_date = pd.to_datetime(end_date)

# Create a new dataframe with all dates between start_date and end_date
all_dates = pd.date_range(start=start_date, end=end_date, freq='D')
all_dates_df = pd.DataFrame({'EpisodeDate': all_dates})

# Merge all_dates_df with datewise_df
datewise_df = pd.merge(all_dates_df, datewise_df, on='EpisodeDate', how='outer').fillna(0)

# Apply the same multiplication to the 'Low Score', 'Good Score', and 'High Score' columns
datewise_df['Low Score'] = datewise_df['Low Score'] * scaling_factor
datewise_df['Good Score'] = datewise_df['Good Score'] * scaling_factor
datewise_df['High Score'] = datewise_df['High Score'] * scaling_factor

no_videos_df = datewise_df[datewise_df['Hand Washes Observed'] == 0]
date_strings = no_videos_df['EpisodeDate'].dt.strftime('%Y-%m-%d').tolist()

videos_df = datewise_df[datewise_df['Hand Washes Observed'] != 0].copy()
videos_df['EpisodeDate'] = pd.to_datetime(videos_df['EpisodeDate']).apply(lambda x: x.strftime('%d %b'))
videos_list = videos_df[['EpisodeDate', 'Hand Washes Observed', 'Low Score', 'Good Score', 'High Score']].values.tolist()

# Print the resulting list of lists
print(videos_list)
#pdb.set_trace()
## Initialize an empty list to store the date strings
#date_list = []

## Iterate over the days in the range and add the corresponding date string to the list
#for day in range(int(start_day), int(end_day) + 1):
#    date_string = f'{day:02d} {start_month_name[:3]}'
#    date_list.append(date_string)

## Print the resulting list of date strings
#print(date_list)



# Print the resulting dataframe
print(datewise_df[['EpisodeDate', 'Hand Washes Observed', 'Low Score', 'Good Score', 'High Score']])

# Calculate counts and merge with step counts
counts_df = df[df['Hand Washes Observed'] == 1]

counts_df = counts_df['core_steps_count'].value_counts().sort_index().reset_index()
print(counts_df)
counts_df.columns = ['Count Out Of Six', 'Number of HandWashes']
print(counts_df)

step_counts_df = pd.DataFrame({'Count Out Of Six': range(1, 7)})
print(step_counts_df)
counts_df = pd.merge(step_counts_df, counts_df, how='left', on='Count Out Of Six').fillna(0)

counts_df['Number of Handwashes'] = counts_df['Number of HandWashes']
total_handwashes = df['Hand Washes Observed'].sum()
# total_handwash = counts_df['Number of HandWashes'].sum()
# pdb.set_trace()
# Determine most common step
most_common_score = counts_df.loc[counts_df['Number of HandWashes'].idxmax(), 'Count Out Of Six']
print(most_common_score)


# Create text labels
# total_handwashes = df['Hand Washes Observed'].sum()

# Calculate percentage values
counts_df['Handwash Percentages'] = (
        counts_df['Number of HandWashes'] / counts_df['Number of HandWashes'].sum()).fillna(0)  # * 100)

#counts_df['text'] = counts_df['Handwash Percentages'].apply(lambda x: f'<b>{x}%</b>')

# Convert dataframe to dictionary
count_dict = counts_df.to_dict(orient='list')
print(count_dict)
text_colors = ['black' if y in [3, 4] else 'black' for y in count_dict['Count Out Of Six']]

# Create trace object
trace1 = go.Bar(
    x=[(count_dict['Number of HandWashes'][i]* scaling_factor) if count_dict['Count Out Of Six'][i] != 0 else 0 for i in
       range(len(count_dict['Number of HandWashes']))],
    y=[count_dict['Count Out Of Six'][i] if count_dict['Count Out Of Six'][i] != 0 else '' for i in
       range(len(count_dict['Count Out Of Six']))],
    marker=dict(
        color=['rgb(153, 0, 0)' if y in [1, 2] else 'rgb(241, 194, 50)' if y in [3, 4] else 'rgb(31, 146, 70)' for y in
               count_dict['Count Out Of Six']]),
    text=[f'<b>{int(round(x* scaling_factor))}</b>{" " * 1}' for x in count_dict['Number of HandWashes']],
    textfont=dict(color=text_colors, size=20, family='Poppins'),
    textposition='auto',
    showlegend=False,
    orientation='h',
    width=[0.5, 0.5, 0.5, 0.5, 0.5, 0.5]
    # offset = 0.0002
)

# round(x*100,1)

# Modify the y-axis labels
y_labels = [str(count) + '/ 6' if count != 0 else '' for count in count_dict['Count Out Of Six']]


# Create trace object for the line
trace2 = go.Scatter(
    x=[0, counts_df['Number of HandWashes'].sum()],  # Set x values for the line (from 0 to rounded_total_handwashes)
    y=[0.5, 0.5],  # Set y values for the line (a horizontal line at y=3.5)
    mode='lines',
    line=dict(color='rgb(103, 70, 203)', width=3),  # Set line color and width
    name='Target',  # Set the name of the trace
    showlegend=False
)

# Combine both traces in one figure
fig = go.Figure(data=[trace1, trace2])

# fig.add_annotation(x=1 / 2, y=0.4,
#                   text=f"TOTAL NUMBER OF HANDWASHES: {int(round(total_handwashes*scaling_factor,0))}", font=dict(color='rgb(103, 70, 203)', size=16, family='Poppins'),
#                   showarrow=False)

print(Device_to_location[DeviceUUID], int(round(total_handwashes * scaling_factor, 0)))

fig.update_layout(
    yaxis=dict(
        linewidth=2,  # Set x-axis line width to 2
        linecolor='black',
        tickangle=-90,
        ticktext=y_labels,
        tickvals=count_dict['Count Out Of Six'],
        tickfont=dict(
            family='Poppins',
            size=20,
            color='black',
        )),
    xaxis=dict(
        # range=[0, rounded_total_handwashes],  # Set the range of values on the x-axis to go from 0 to 30
        # type='category',  # Set x-axis to display categorical values
        #tickvals=[0.2, 0.4, 0.6, 0.8, 1],
        #tickformat='.0%',
        linewidth=2,
        linecolor='black',
        tickfont=dict(
            family='Poppins',
            size=20,
            color='black',
        ),
        showgrid=True,  # Set showgrid to True to show grid lines on x-axis
        gridwidth=1,  # Set the width of grid lines on x-axis
        gridcolor='lightgray'  # Set the color of grid lines on x-axi
    ),
    plot_bgcolor='white',  # Set background color to white
    bargap=0.2,
    bargroupgap=0.1,
    height=800,  # set the height to 800 pixels
    width=600,  # set the width to 600 pixels
    margin=dict(l=20, r=20, t=20, b=20)
)

# Display the graph
# Determine filename and save image
filename = 'stepcount_bar_chart.png'
fig.write_image(filename, engine='kaleido')

# Sample data
core_steps = ['PalmtoPalm_Detected', 'PalmtoDorsum_Detected', 'FingersInterlaced_Detected', 'FistInterlocked_Detected',
              'ThumbRub_Detected', 'PalmtoNails_Detected']

output = {}
common_score = 0
for step in core_steps:
    step_compliance = (df[step].sum(axis=0) / (
                df['Hand Washes Observed'].sum(axis=0) * scaling_factor))  # * f_factor[DeviceUUID]))
    if numpy.isnan(step_compliance):
        output[step] = 0
    elif step_compliance > 1:
        output[step] = 1
    else:
        output[step] = step_compliance

step_df = pd.DataFrame(output, index=[0])
step_df = step_df.transpose().reset_index()
step_df.rename({'index': 'step_name', 0: 'compliance_score'}, inplace=True, axis=1)
step_df['step_name'] = step_df['step_name'].replace(
    {'PalmtoPalm_Detected': 'Palm to Palm', 'PalmtoDorsum_Detected': 'Palm to Dorsum',
     'FingersInterlaced_Detected': 'Fingers Interlaced', 'FistInterlocked_Detected': 'Fist Interlocked',
     'ThumbRub_Detected': 'Thumb Rub', 'PalmtoNails_Detected': 'Palm to Nails'}, regex=True)
step_df['Color'] = step_df['compliance_score'].apply(lambda
                                                         x: '<b>Less than 50% Compliance<br>(Average)</b>' if x <= 0.5 else '<b>More than 50% Compliance<br>(Average)</b>')
step_df['Text'] = step_df['compliance_score'].apply(lambda x: "<b>" + str(round(x * 100)) + "</b>%")

print(step_df)
# Convert dataframe to dictionary
data_dict = step_df.to_dict(orient='list')

trace1 = go.Bar(
    x=data_dict['step_name'],
    y=[y * 100 for y in data_dict['compliance_score']],  # Multiply by 100 to convert to percentage format
    marker=dict(color=['rgb(31, 146, 70)' if y >= 0.5 else 'rgb(153, 0, 0)' for y in data_dict['compliance_score']]),
    width=[0.4, 0.4, 0.4, 0.4, 0.4, 0.4],
    name='Percentage',
    # text=[f'{" " * 5}<b>{x}</b>' for x in data_dict['step_name']],

    showlegend=False,
    # offset=0.2
)

# Create trace for annotations
# annotations = []
# for x, y in zip(data_dict['step_name'], data_dict['compliance_score']):
#    annotations.append(dict(x=x, y=5, text=f'<b>{x}</b>', font=dict(size=20, color='white', family='assets\Poppins\Poppins-Regular.tff'), showarrow=False,
#                            textangle=-90, valign='bottom', xanchor='center', yanchor='bottom'))

# Create trace for translucent bar
trace2 = go.Bar(
    x=data_dict['step_name'],
    y=[(1 - y) * 100 for y in data_dict['compliance_score']],
    # Subtract compliance score from 1 and multiply by 100 to get the non-compliant percentage
    marker=dict(
        color=['rgb(31, 146, 70)' if y >= 0.5 else 'rgb(153, 0, 0)' for y in data_dict['compliance_score']],
        opacity=0.2),
    width=[0.4, 0.4, 0.4, 0.4, 0.4, 0.4],
    showlegend=False,
    name='Rest',
    # offset=0.2
)

# Combine both traces in one figure
fig = go.Figure(data=[trace1, trace2])  # , layout=go.Layout(annotations=annotations))

# Update layout options
fig.update_layout(
    yaxis=dict(  # tickformat='%',
        linewidth=2,  # Set x-axis line width to 2
        # tickvals=[0, 25, 50, 75, 100],
        ticktext=['<b>0%</b>', '<b>25%</b>', '<b>50%</b>', '<b>75%</b>', '<b>100%</b>'],
        tickangle=-90,
        linecolor='black',
        tickfont=dict(size=20, family='Poppins')),
    xaxis=dict(
        categoryorder='array',
        categoryarray=data_dict['step_name'],
        linewidth=2,
        linecolor='black'
    ),
    title=None,
    xaxis_title=None,
    yaxis_title=None,
    bargap=0.2,
    bargroupgap=0.1,
    plot_bgcolor='white',
    height=800,  # set the height to 800 pixels
    width=600,  # set the width to 600 pixels
    margin=dict(l=20, r=20, t=20, b=20)
)

# Set the y-axis range to start from 0 and end at 100%
fig.update_yaxes(range=[0, 100])

# Set the mode of the bar chart to "stack" to stack the bars on top of each other
fig.update_layout(barmode='stack')

# Hide x-axis labels
fig.update_xaxes(showticklabels=False)

# Save the figure as a PNG file
pio.write_image(fig, 'percentage_bar_graph.png')

# -----------------------------------------------------------


# Access the presentation template
if site_type == 'Room':
    pr_ppt = Presentation('Performance report - LR Template with Data Table - English - NS.pptx') #'Performance report - LR Template - English - NS.pptx')
else:
    pr_ppt = Presentation('Performance report - SL Template With Extra Rows- English - NS.pptx')

# Changes to Slide 1
text_into_presentation(pr_ppt, 0, 5, 4, 0, Device_to_location[DeviceUUID])
text_into_presentation(pr_ppt, 0, 5, 4, 1, date_range_text)

# Changes to Slide 2
text_into_presentation(pr_ppt, 1, 2, 0, 1, date_range_text)
text_into_presentation(pr_ppt, 1, 2, 5, 0, str(most_common_score))  # common_step + " Out Of 6*")
img_into_presentation(pr_ppt, 1, "stepcount_bar_chart.png", 3, 4.7, 5, 5.3)

added_text = ''

print(most_common_score)
#print(prev_common_score)
if prev_category == highest_category:
    added_text = '\n\n Same as last time.'
elif prev_category>highest_category:  #prev_common_score in [3,4] and most_common_score in [1,2]:
    added_text = '\n\n Worse than last time.'
elif prev_category<highest_category: #(prev_common_score in [1,2] and most_common_score in [3,4]) or (prev_common_score in [3,4] and most_common_score in [5,6]):
    added_text = '\n\n Better than last time.'

low_score = highest_category in [0]
good_score = highest_category in [1]
high_score = highest_category in [2]

slide_no = 1  # the index of the slide containing the existing shape
shape_no = 4  # the shape number of the existing shape
if high_score:
    text = 'Excellent!\nSupport your team to sustain their top performance.' + added_text
else:
    text = 'Excellent!\nSupport your team to sustain their top performance.'
color = (31, 146, 70)  # the RGB color values as a tuple
left = 1000000  # 902299
top = 4768750
width = 1594800
height = 585000
print(text)
print(high_score)
add_text_to_shape(pr_ppt, slide_no, shape_no, text, color, high_score, left, top, width, height)

shape_no = 5  # he shape number of the existing shape
if good_score:
    text = 'Well Done!\nAcknowledge your team’s effort and support them to improve.' + added_text
else:
    text = 'Well Done!\nAcknowledge your team’s effort and support them to improve.'
color = (241, 194, 50)  # the RGB color values as a tuple
left = 1000000  # 902299
top = 6000000  # 7306525
width = 1594800
height = 875030
print(text)
print(good_score)
add_text_to_shape(pr_ppt, slide_no, shape_no, text, color, good_score, left, top, width, height)

shape_no = 6  # he shape number of the existing shape
if low_score:
    text = 'Can Do Better.\nAsk your team, what will help them learn and follow the steps better.' + added_text
else:
    text = 'Can Do Better.\nAsk your team, what will help them learn and follow the steps better.'
color = (153, 0, 0)  # the RGB color values as a tuple
left = 1000000  # 902299
top = 7500000  # 8889088
width = 1594800
height = 400200
print(text)
print(low_score)
add_text_to_shape(pr_ppt, slide_no, shape_no, text, color, low_score, left, top, width, height)

text = f'TOTAL NO. OF HANDWASHES: {round(total_handwashes * scaling_factor)}'
# + added_text
color = (103, 70, 203)  # the RGB color values as a tuple
left = Inches(4.597916666666666)  # 4.597916666666666 9.4875 2.4256944444444444 0.23402777777777778
top = Inches(9.4875)
width = Inches(2.4256944444444444)
height = Inches(0.23402777777777778)
print(text)
print(good_score)
add_text_to_shape2(pr_ppt, slide_no, text, color, left, top, width, height, 0)

# Changes to Slide 3
img_into_presentation(pr_ppt, 2, "percentage_bar_graph.png", 0.4, 1.9, 4.3, 4.5)

step_labels = [
    ['Palm to Palm', 0.8284722222222223, 4.7131944444444445, 0.35138888888888886, 1.2798611111111111],
    ['Palm to Dorsum', 1.4673611111111111, 4.71259842519685, 0.35138888888888886, 1.2798611111111111],
    ['Fingers Interlaced', 2.13125, 4.71259842519685, 0.35138888888888886, 1.2798611111111111],
    ['Fists Interlocked', 2.763888888888889, 4.71259842519685, 0.35138888888888886, 1.2798611111111111],
    ['Thumb Rub', 3.4277777777777776, 4.71259842519685, 0.35138888888888886, 1.2798611111111111],
    ['Nails to Palm', 4.060416666666667, 4.71259842519685, 0.35138888888888886, 1.2798611111111111]
]

for values in step_labels:
    # + added_text
    print(values)
    text = values[0]
    color = (255, 255, 255)  # the RGB color values as a tuple
    left = Inches(values[1] + 0.52)  # 902299
    top = Inches(values[2])  # 7306525
    width = Inches(values[3])
    height = Inches(values[4])
    print(text)
    print(good_score)
    add_text_to_shape2(pr_ppt, 2, text, color, left, top, width, height, -90)

if site_type == 'Room':

    for i in range(2, 8):
        pr_ppt.slides[2].shapes[i].left = Inches(0.8 + ((i - 2) * 0.65))
        pr_ppt.slides[2].shapes[i].top = Inches(6.35)

    # Changes to Slide 4
    for i, date in enumerate(date_strings):
        table_cell_into_presentation(pr_ppt, 3, 3, i + 1, 0, date)

    pr_ppt.save(
        'Performance report -' + Device_to_location[DeviceUUID] + ' - ' + date_title_text + '- English.pptx')

else:
    for i in range(0, 6):
        pr_ppt.slides[2].shapes[i].left = Inches(0.9 + (i * 0.65))
        pr_ppt.slides[2].shapes[i].top = Inches(6.35)

    # Changes to Slide 4
    text_into_presentation(pr_ppt, 3, 1, 1, 0, 'Period under review: ' + date_range_text)

    for i, video_data in enumerate(videos_list):
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 0, video_data[0])
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 1, str(round(video_data[1])))
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 2, str(round(video_data[2])))
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 3, str(round(video_data[3])))
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 4, str(round(video_data[4])))

    # get the first slide
    slide = pr_ppt.slides[3]

    # get the first table on the slide
    table = slide.shapes[2].table

    # Create a list to keep track of the rows with the highest values in each column
    max_rows = [None] * (len(table.columns) - 1)

    # Loop through the columns in the table
    for col_idx, column in enumerate(table.columns):
        # Skip the first column
        if col_idx == 0:
            continue

        # Set the starting row to 1 to skip the header row
        start_row = 1

        # Set the starting value to 0
        max_value = 0

        # Select the cells in the column and loop through them
        for row_idx, cell in enumerate([row.cells[col_idx] for row in list(table.rows)[start_row:]]):
            # Get the value of the cell as a float
            #value = float(cell.text)
            try:
                value = float(cell.text)
            except ValueError:
                value = 0

            # If the value is greater than the current max value, update the max value and max cell
            if value > max_value:
                max_value = value
                max_cell = cell
                max_rows[col_idx - 1] = row_idx + start_row

    max_rows_dict = {}
    for col_idx, row_idx in enumerate(max_rows):
        if row_idx is not None:
            if row_idx not in max_rows_dict:
                max_rows_dict[row_idx] = []
            max_rows_dict[row_idx].append(col_idx + 1)

    print(max_rows_dict)

    # Loop through the rows with the highest values in each column
    for row_idx in max_rows_dict.keys():
        row_max_value = 0
        row_max_cell = []
        for col_idx in max_rows_dict[row_idx]:
            if col_idx < 2:
                continue
            cell = table.cell(row_idx, col_idx)

            # Get the value of the cell as a float
            try:
                value = float(cell.text)
            except ValueError:
                value = 0

            # Check if this is the highest value in the row so far
            if value > row_max_value:
                row_max_value = value
                row_max_column = [col_idx]
                row_max_cell = [cell]
            elif value == row_max_value:
                row_max_column.append(col_idx)
                row_max_cell.append(cell)

        # assign the color based on the column index
        for cell in row_max_cell:
            if row_max_value == 0:
                color = None  # default color if no row max value found
            elif row_max_value == float(cell.text):
                if row_max_column[0] == 2:
                    color = RGBColor(252, 192, 194)
                elif row_max_column[0] == 3:
                    color = RGBColor(241, 194, 50)
                elif row_max_column[0] == 4:
                    color = RGBColor(255, 242, 204)
            else:
                color = None  # default color for non-maximum cells
            if color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])

    pr_ppt.save(
        'Performance report -' + Device_to_location[DeviceUUID] + ' - ' + date_title_text + '- English.pptx')
