import calendar
import json
import math
import numpy as np

import pandas as pd
import plotly.graph_objs as go
import plotly.io as pio
import requests

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

import datetime
import argparse

# Local imports (sanitized config/helpers)
from config import Device_to_location, compliance_data_link
from helpers import (
    text_into_presentation,
    table_cell_into_presentation,
    img_into_presentation,
    add_text_to_shape,
    add_text_to_shape2,
)

import pdb

# -------------------------------
# Argument parsing
# -------------------------------
parser = argparse.ArgumentParser()
parser.add_argument('--uuid', default="SAMPLE_UUID")
parser.add_argument('--date_start', default='2023-04-01')
parser.add_argument('--date_end', default='2023-04-30')
parser.add_argument('--time_start', default='00:00:01')
parser.add_argument('--time_end', default='23:59:59')
parser.add_argument('--prev_start', default='2023-03-01')
parser.add_argument('--prev_end', default='2023-03-31')
args = parser.parse_args()

# -------------------------------
# Constants / inputs
# -------------------------------
print("Setting up constants")

DeviceUUID = args.uuid
start_date = args.date_start
end_date = args.date_end

start_time, end_time = args.time_start, args.time_end
prev_start, prev_end = args.prev_start, args.prev_end

# Date formatting for titles
start_year, start_month, start_day = start_date.split('-')
_, end_month, end_day = end_date.split('-')

start_month_name = calendar.month_name[int(start_month)]
end_month_name = calendar.month_name[int(end_month)]

date_range_text = f'{start_day} to {end_day} {start_month_name[:3]} {start_year}'
date_title_text = f'{end_day} {start_month_name[:3]} {start_year}'
print(date_range_text)
print(date_title_text)

# Load scaling factors (file name kept for backward compatibility)
# For OSS repo: include a sample_scaling_factors.json and add real one to .gitignore
with open('f_factor.json', 'r') as a:
    f_factor = json.load(a)

print("Loaded scaling factors (sample or real):", f_factor)

# -------------------------------
# Helper: safe float cast
# -------------------------------
def safe_float_col(df, col, default=0.0):
    if col in df.columns:
        try:
            df[col] = df[col].astype(float)
        except Exception:
            # attempt elementwise conversion
            df[col] = pd.to_numeric(df[col], errors='coerce').fillna(default)
    else:
        df[col] = default
    return df

# -------------------------------
# Retrieve data for current period
# -------------------------------
url = f"{compliance_data_link}/get_specific_compliance"
params = {'installations': DeviceUUID, 'start_date': start_date, 'end_date': end_date}
response = requests.get(url, params=params)
d = json.loads(response.text)
df = pd.DataFrame(d)

# Normalise booleans and set recorded flag
df.replace({'True': 1, 'False': 0}, inplace=True)
df['Recorded Events'] = 1

# Try to normalise a generic duration column if present (keeps compatibility)
if 'TotalWashDuration' in df.columns:
    df = safe_float_col(df, 'TotalWashDuration', default=0.0)
elif 'TotalDuration' in df.columns:
    df = safe_float_col(df, 'TotalDuration', default=0.0)
# else: no duration column available; that's fine

# Dynamically discover feature columns (previously handwash steps)
feature_cols = [c for c in df.columns if c.endswith('_Detected')]
num_features = max(1, len(feature_cols))  # keep >=1 to avoid division by zero later
print("Detected feature columns:", feature_cols)

# Create generic counts and event categories
df['core_features_count'] = df[feature_cols].sum(axis=1) if feature_cols else 0
df['Event Observed'] = df['core_features_count'].apply(lambda x: 1 if x > 0 else 0)
df['Low Score'] = df['core_features_count'].apply(lambda x: 1 if x in [1, 2] else 0)
df['Good/Medium Score'] = df['core_features_count'].apply(lambda x: 1 if x in [3, 4] else 0)
df['High Score'] = df['core_features_count'].apply(lambda x: 1 if x in [5, 6] else 0)
df['Full Score'] = df['core_features_count'].apply(lambda x: 1 if x == num_features else 0)

# Scaling factor safely retrieved
scaling_factor = f_factor.get(DeviceUUID, f_factor.get('Median', 1))
print("Using scaling factor:", scaling_factor)

# Parse and normalise date/time columns
if 'EpisodeTime' in df.columns:
    # original format expected integer like HHMMSS stored as int/str; keep original parsing
    def parse_time(x):
        try:
            x = int(x)
            hh = int(x / 10000)
            mm = int((x // 100) % 100)
            ss = int(x % 100)
            return f"{hh:02d}:{mm:02d}:{ss:02d}"
        except Exception:
            return "00:00:00"
    df['EpisodeTime'] = df['EpisodeTime'].apply(parse_time)
else:
    df['EpisodeTime'] = "00:00:00"

df["DateTime"] = pd.to_datetime(df['EpisodeDate'].astype(str) + " " + df['EpisodeTime'].astype(str), errors='coerce')
df['EpisodeDate'] = pd.to_datetime(df['EpisodeDate'], errors='coerce')

# -------------------------------
# Aggregate counts & categories
# -------------------------------
low_counts_df = int(df['Low Score'].sum())
good_counts_df = int(df['Good/Medium Score'].sum())
high_counts_df = int(df['High Score'].sum())

max_counts = max(low_counts_df, good_counts_df, high_counts_df)
if max_counts == high_counts_df:
    highest_category = 2
elif max_counts == good_counts_df:
    highest_category = 1
else:
    highest_category = 0

# -------------------------------
# Retrieve previous period (if needed) and compute comparison category
# -------------------------------
prev_params = {'installations': DeviceUUID, 'start_date': prev_start, 'end_date': prev_end}
response = requests.get(url, params=prev_params)
d_prev = json.loads(response.text)
prev_df = pd.DataFrame(d_prev)

prev_df.replace({'True': 1, 'False': 0}, inplace=True)
prev_df['Recorded Events'] = 1
prev_df = safe_float_col(prev_df, 'TotalWashDuration', default=0.0)

prev_feature_cols = [c for c in prev_df.columns if c.endswith('_Detected')]
prev_df['core_features_count'] = prev_df[prev_feature_cols].sum(axis=1) if prev_feature_cols else 0
prev_df['Low Score'] = prev_df['core_features_count'].apply(lambda x: 1 if x in [1, 2] else 0)
prev_df['Good/Medium Score'] = prev_df['core_features_count'].apply(lambda x: 1 if x in [3, 4] else 0)
prev_df['High Score'] = prev_df['core_features_count'].apply(lambda x: 1 if x in [5, 6] else 0)

low_counts_prev_df = int(prev_df['Low Score'].sum())
good_counts_prev_df = int(prev_df['Good/Medium Score'].sum())
high_counts_prev_df = int(prev_df['High Score'].sum())

max_counts_prev = max(low_counts_prev_df, good_counts_prev_df, high_counts_prev_df)
if max_counts_prev == high_counts_prev_df:
    prev_category = 2
elif max_counts_prev == good_counts_prev_df:
    prev_category = 1
else:
    prev_category = 0

# -------------------------------
# Datewise aggregation (fill missing dates)
# -------------------------------
site_type = Device_to_location.get(DeviceUUID, "Unknown Location").split(" ")[-1]

datewise_df = df.groupby(['EpisodeDate']).sum(numeric_only=True).reset_index()
# Apply scaling to counts
for col in ['Event Observed', 'Low Score', 'Good/Medium Score', 'High Score']:
    if col in datewise_df.columns:
        datewise_df[col] = datewise_df[col] * scaling_factor
# Ensure full date range present
start_dt = pd.to_datetime(start_date)
end_dt = pd.to_datetime(end_date)
all_dates = pd.date_range(start=start_dt, end=end_dt, freq='D')
all_dates_df = pd.DataFrame({'EpisodeDate': all_dates})
datewise_df = pd.merge(all_dates_df, datewise_df, on='EpisodeDate', how='outer').fillna(0)

no_event_days = datewise_df[datewise_df.get('Event Observed', 0) == 0]
date_strings = no_event_days['EpisodeDate'].dt.strftime('%Y-%m-%d').tolist()

videos_df = datewise_df[datewise_df.get('Event Observed', 0) != 0].copy()
videos_df['EpisodeDate'] = pd.to_datetime(videos_df['EpisodeDate']).apply(lambda x: x.strftime('%d %b'))
videos_list = videos_df[['EpisodeDate', 'Event Observed', 'Low Score', 'Good/Medium Score', 'High Score']].fillna(0).values.tolist()

print("Daily non-zero entries (sample):", videos_list[:5])

# -------------------------------
# Distribution of core feature counts
# -------------------------------
counts_df = df[df['Event Observed'] == 1]
counts_df = counts_df['core_features_count'].value_counts().sort_index().reset_index()
counts_df.columns = ['Count Out Of N', 'Number of Events']

# Ensure we have rows for full range 1..num_features
step_counts_df = pd.DataFrame({'Count Out Of N': list(range(1, num_features + 1))})
counts_df = pd.merge(step_counts_df, counts_df, how='left', on='Count Out Of N').fillna(0)

total_events = int(df['Event Observed'].sum())

most_common_score = int(counts_df.loc[counts_df['Number of Events'].idxmax(), 'Count Out Of N']) if counts_df['Number of Events'].sum() > 0 else 0
print("Most common count:", most_common_score)

# Percentages (safe)
total_num_events = counts_df['Number of Events'].sum()
if total_num_events == 0:
    counts_df['Percent'] = 0.0
else:
    counts_df['Percent'] = counts_df['Number of Events'] / total_num_events

count_dict = counts_df.to_dict(orient='list')
print("Count dictionary:", count_dict)

# -------------------------------
# Plot 1: horizontal distribution bar
# -------------------------------
trace1 = go.Bar(
    x=[(count_dict['Number of Events'][i] * scaling_factor) if count_dict['Count Out Of N'][i] != 0 else 0 for i in range(len(count_dict['Number of Events']))],
    y=[count_dict['Count Out Of N'][i] if count_dict['Count Out Of N'][i] != 0 else '' for i in range(len(count_dict['Count Out Of N']))],
    marker=dict(
        color=['rgb(153, 0, 0)' if y in [1, 2] else 'rgb(241, 194, 50)' if y in [3, 4] else 'rgb(31, 146, 70)' for y in count_dict['Count Out Of N']]),
    text=[f'<b>{int(round(x * scaling_factor))}</b>' for x in count_dict['Number of Events']],
    textfont=dict(color='black', size=16, family='Poppins'),
    textposition='auto',
    showlegend=False,
    orientation='h',
    width=[0.5] * len(count_dict['Count Out Of N'])
)

y_labels = [f"{int(count)}/ {num_features}" if count != 0 else '' for count in count_dict['Count Out Of N']]

trace2 = go.Scatter(
    x=[0, counts_df['Number of Events'].sum()],
    y=[0.5, 0.5],
    mode='lines',
    line=dict(color='rgb(103, 70, 203)', width=3),
    showlegend=False
)

fig = go.Figure(data=[trace1, trace2])

fig.update_layout(
    yaxis=dict(
        linewidth=2,
        linecolor='black',
        tickangle=-90,
        ticktext=y_labels,
        tickvals=count_dict['Count Out Of N'] if count_dict['Count Out Of N'] else [],
        tickfont=dict(family='Poppins', size=14, color='black'),
    ),
    xaxis=dict(
        linewidth=2,
        linecolor='black',
        tickfont=dict(family='Poppins', size=14, color='black'),
        showgrid=True,
        gridwidth=1,
        gridcolor='lightgray'
    ),
    plot_bgcolor='white',
    bargap=0.2,
    bargroupgap=0.1,
    height=800,
    width=600,
    margin=dict(l=20, r=20, t=20, b=20)
)

filename = 'output/feature_count_distribution.png'
fig.write_image(filename, engine='kaleido')
print("Saved distribution chart to", filename)

# -------------------------------
# Per-feature compliance / coverage chart
# -------------------------------
# Build step/feature compliance from feature_cols
output = {}
for step in feature_cols:
    denom = max(1, int(df['Event Observed'].sum() * scaling_factor))
    step_sum = int(df[step].sum()) if step in df.columns else 0
    step_compliance = step_sum / denom if denom > 0 else 0
    # clamp 0..1
    step_compliance = max(0.0, min(1.0, step_compliance))
    output[step] = step_compliance

if feature_cols:
    step_df = pd.DataFrame(output, index=[0]).transpose().reset_index()
    step_df.rename({'index': 'feature_name', 0: 'coverage_score'}, inplace=True, axis=1)
    # nice readable names: Feature 1, Feature 2, ...
    readable_names = {feature_cols[i]: f'Feature {i+1}' for i in range(len(feature_cols))}
    step_df['feature_name_readable'] = step_df['feature_name'].replace(readable_names)
    step_df['Label'] = step_df['coverage_score'].apply(lambda x: '<b>Less than 50% Coverage</b>' if x <= 0.5 else '<b>More than 50% Coverage</b>')
    step_df['Text'] = step_df['coverage_score'].apply(lambda x: "<b>" + str(round(x * 100)) + "</b>%")
else:
    step_df = pd.DataFrame({
        'feature_name': [],
        'coverage_score': [],
        'feature_name_readable': [],
        'Label': [],
        'Text': []
    })

print(step_df)

data_dict = {
    'feature_name': step_df['feature_name_readable'].tolist(),
    'coverage_score': step_df['coverage_score'].tolist()
}

trace_a = go.Bar(
    x=data_dict['feature_name'],
    y=[y * 100 for y in data_dict['coverage_score']],
    marker=dict(color=['rgb(31, 146, 70)' if y >= 0.5 else 'rgb(153, 0, 0)' for y in data_dict['coverage_score']]),
    width=[0.4] * len(data_dict['feature_name']),
    showlegend=False,
)

trace_b = go.Bar(
    x=data_dict['feature_name'],
    y=[(1 - y) * 100 for y in data_dict['coverage_score']],
    marker=dict(
        color=['rgb(31, 146, 70)' if y >= 0.5 else 'rgb(153, 0, 0)' for y in data_dict['coverage_score']],
        opacity=0.2),
    width=[0.4] * len(data_dict['feature_name']),
    showlegend=False
)

fig2 = go.Figure(data=[trace_a, trace_b])

fig2.update_layout(
    yaxis=dict(
        linewidth=2,
        ticktext=['<b>0%</b>', '<b>25%</b>', '<b>50%</b>', '<b>75%</b>', '<b>100%</b>'],
        tickangle=-90,
        linecolor='black',
        tickfont=dict(size=14, family='Poppins')
    ),
    xaxis=dict(categoryorder='array', categoryarray=data_dict['feature_name'], linewidth=2, linecolor='black'),
    plot_bgcolor='white',
    height=800,
    width=600,
    margin=dict(l=20, r=20, t=20, b=20),
    bargap=0.2,
    bargroupgap=0.1
)

fig2.update_yaxes(range=[0, 100])
fig2.update_layout(barmode='stack')
fig2.update_xaxes(showticklabels=False)

pio.write_image(fig2, 'output/feature_coverage_chart.png')
print("Saved coverage chart to output/feature_coverage_chart.png")

# -------------------------------
# Load templates and write PPTX slides (generic templates)
# -------------------------------
if site_type == 'Room':
    pr_ppt = Presentation('templates/lr_template_generic.pptx')
else:
    pr_ppt = Presentation('templates/sl_template_generic.pptx')

# Slide 1: Title & date
text_into_presentation(pr_ppt, 0, 5, 4, 0, Device_to_location.get(DeviceUUID, "Unknown Location"))
text_into_presentation(pr_ppt, 0, 5, 4, 1, date_range_text)

# Slide 2: distribution chart and summary
text_into_presentation(pr_ppt, 1, 2, 0, 1, date_range_text)
text_into_presentation(pr_ppt, 1, 2, 5, 0, str(most_common_score))
img_into_presentation(pr_ppt, 1, "output/feature_count_distribution.png", 3, 4.7, 5, 5.3)

# Add assessment boxes (generic text)
added_text = ''
if prev_category == highest_category:
    added_text = '\n\n Same as last period.'
elif prev_category > highest_category:
    added_text = '\n\n Worse than last period.'
elif prev_category < highest_category:
    added_text = '\n\n Better than last period.'

low_score_flag = highest_category == 0
good_score_flag = highest_category == 1
high_score_flag = highest_category == 2

slide_no = 1
shape_no = 4
if high_score_flag:
    text = 'Excellent!\nSupport the team to sustain top performance.' + added_text
else:
    text = 'Excellent!\nSupport the team to sustain top performance.'
color = (31, 146, 70)
left = 1000000
top = 4768750
width = 1594800
height = 585000
add_text_to_shape(pr_ppt, slide_no, shape_no, text, color, high_score_flag, left, top, width, height)

shape_no = 5
if good_score_flag:
    text = 'Well Done!\nAcknowledge the effort and support improvement.' + added_text
else:
    text = 'Well Done!\nAcknowledge the effort and support improvement.'
color = (241, 194, 50)
left = 1000000
top = 6000000
width = 1594800
height = 875030
add_text_to_shape(pr_ppt, slide_no, shape_no, text, color, good_score_flag, left, top, width, height)

shape_no = 6
if low_score_flag:
    text = 'Can Do Better.\nAsk the team what will help improve performance.' + added_text
else:
    text = 'Can Do Better.\nAsk the team what will help improve performance.'
color = (153, 0, 0)
left = 1000000
top = 7500000
width = 1594800
height = 400200
add_text_to_shape(pr_ppt, slide_no, shape_no, text, color, low_score_flag, left, top, width, height)

# Total events label (generic)
text = f'TOTAL NO. OF EVENTS: {round(total_events * scaling_factor)}'
color = (103, 70, 203)
left = Inches(4.597916666666666)
top = Inches(9.4875)
width = Inches(2.4256944444444444)
height = Inches(0.23402777777777778)
add_text_to_shape2(pr_ppt, slide_no, text, color, left, top, width, height, 0)

# Slide 3: coverage chart
img_into_presentation(pr_ppt, 2, "output/feature_coverage_chart.png", 0.4, 1.9, 4.3, 4.5)

# Put generic feature labels around the chart if features exist
feature_labels = [f'Feature {i+1}' for i in range(len(feature_cols))]
for i, label in enumerate(feature_labels):
    left = Inches(0.9 + (i * 0.65))
    top = Inches(6.35)
    width = Inches(0.35)
    height = Inches(1.28)
    add_text_to_shape2(pr_ppt, 2, label, (255, 255, 255), left, top, width, height, -90)

# Slide 4: table of dates without events or with counts
if site_type == 'Room':
    # small layout tweaks for Room template (keeps same behaviour but generic)
    for i in range(2, 8):
        try:
            pr_ppt.slides[2].shapes[i].left = Inches(0.8 + ((i - 2) * 0.65))
            pr_ppt.slides[2].shapes[i].top = Inches(6.35)
        except Exception:
            pass

    # Fill a column with no-event dates
    for i, date in enumerate(date_strings):
        table_cell_into_presentation(pr_ppt, 3, 3, i + 1, 0, date)

    out_name = f'output/Report-{Device_to_location.get(DeviceUUID,"Unknown")}-{date_title_text}.pptx'
    pr_ppt.save(out_name)
    print("Saved report to", out_name)
else:
    # layout tweaks for other template
    for i in range(0, 6):
        try:
            pr_ppt.slides[2].shapes[i].left = Inches(0.9 + (i * 0.65))
            pr_ppt.slides[2].shapes[i].top = Inches(6.35)
        except Exception:
            pass

    text_into_presentation(pr_ppt, 3, 1, 1, 0, 'Period under review: ' + date_range_text)

    for i, video_data in enumerate(videos_list):
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 0, video_data[0])
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 1, str(round(video_data[1])))
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 2, str(round(video_data[2])))
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 3, str(round(video_data[3])))
        table_cell_into_presentation(pr_ppt, 3, 2, i + 1, 4, str(round(video_data[4])))

    # highlight max values in each row (keeps original logic)
    try:
        slide = pr_ppt.slides[3]
        table = slide.shapes[2].table

        max_rows = [None] * (len(table.columns) - 1)
        for col_idx, column in enumerate(table.columns):
            if col_idx == 0:
                continue
            start_row = 1
            max_value = 0
            for row_idx, cell in enumerate([row.cells[col_idx] for row in list(table.rows)[start_row:]]):
                try:
                    value = float(cell.text)
                except Exception:
                    value = 0
                if value > max_value:
                    max_value = value
                    max_rows[col_idx - 1] = row_idx + start_row

        max_rows_dict = {}
        for col_idx, row_idx in enumerate(max_rows):
            if row_idx is not None:
                max_rows_dict.setdefault(row_idx, []).append(col_idx + 1)

        for row_idx in max_rows_dict.keys():
            row_max_value = 0
            row_max_cell = []
            for col_idx in max_rows_dict[row_idx]:
                if col_idx < 2:
                    continue
                cell = table.cell(row_idx, col_idx)
                try:
                    value = float(cell.text)
                except Exception:
                    value = 0
                if value > row_max_value:
                    row_max_value = value
                    row_max_column = [col_idx]
                    row_max_cell = [cell]
                elif value == row_max_value:
                    row_max_column.append(col_idx)
                    row_max_cell.append(cell)

            for cell in row_max_cell:
                if row_max_value == 0:
                    color = None
                elif row_max_value == float(cell.text):
                    if row_max_column[0] == 2:
                        color = RGBColor(252, 192, 194)
                    elif row_max_column[0] == 3:
                        color = RGBColor(241, 194, 50)
                    elif row_max_column[0] == 4:
                        color = RGBColor(255, 242, 204)
                    else:
                        color = None
                else:
                    color = None
                if color:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])
    except Exception as e:
        print("Table-highlighting step failed (non-fatal):", e)

    out_name = f'output/Report-{Device_to_location.get(DeviceUUID,"Unknown")}-{date_title_text}.pptx'
    pr_ppt.save(out_name)
    print("Saved report to", out_name)
